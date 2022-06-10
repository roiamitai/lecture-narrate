import os
import unittest
from .implementation import slide_extract_lecture_notes
from pptx import Presentation
#
class UnitTest (unittest.TestCase):
    #
    def test_lecture_note_seq (self):
        lecture = Presentation(os.path.join(os.path.dirname(__file__),'test-resources','basic-pres-w-img-example.pptx'))
        self.assertEqual('Notes for title slide.', slide_extract_lecture_notes(lecture.slides[0]))
        self.assertEqual('Notes for second slide.',slide_extract_lecture_notes(lecture.slides[1]))
    #
#